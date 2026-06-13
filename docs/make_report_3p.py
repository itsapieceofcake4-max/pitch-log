# -*- coding: utf-8 -*-
"""
make_report_3p.py
=================
Pitch Log 提案（3枚・レポート中心版）を生成する。
出力: docs/Pitch_Log_Report_3p.pptx
構成: 1=主張（なぜ必要か／因果である） 2=出力レポート（サンプル） 3=因果である根拠
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

BG="0E1825"; PANEL="16273B"; CARD="1A2B41"; CARD2="21364F"; BORDER="2C4A6E"
WHITE="FFFFFF"; INK="E8F0FA"; MUTE="93AAC6"; FAINT="647E9B"
BLUE="3B9EFF"; TEAL="2BD4A8"; RED="FF5C5C"; GOLD="FFC24B"; GREY="8FA0B4"
JP="Meiryo"; EMU_IN=914400


def C(h): return RGBColor.from_string(h)


prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]


def slide(bg=BG):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=C(bg); r.line.fill.background(); r.shadow.inherit=False
    return s


def text(s,t,x,y,w,h,size,color=INK,bold=False,align=PP_ALIGN.LEFT,font=JP,anchor=MSO_ANCHOR.TOP,spacing=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,ln in enumerate(t.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=spacing
        r=p.add_run(); r.text=ln
        f=r.font; f.size=Pt(size); f.bold=bold; f.name=font; f.color.rgb=C(color)
    return tb


def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,radius=0.08,shape=MSO_SHAPE.ROUNDED_RECTANGLE,dash=None):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
    if line is None: sp.line.fill.background()
    else:
        sp.line.color.rgb=C(line); sp.line.width=Pt(lw)
        if dash: sp.line.dash_style=dash
    try:
        if shape==MSO_SHAPE.ROUNDED_RECTANGLE: sp.adjustments[0]=radius
    except Exception: pass
    return sp


def line(s,x1,y1,x2,y2,color=BORDER,w=1.0,dash=None):
    ln=s.shapes.add_connector(2,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    ln.line.color.rgb=C(color); ln.line.width=Pt(w); ln.shadow.inherit=False
    if dash: ln.line.dash_style=dash
    return ln


def oval(s,cx,cy,r,fill,lineC=None,lw=1.0):
    sp=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cx-r),Inches(cy-r),Inches(2*r),Inches(2*r)); sp.shadow.inherit=False
    sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
    if lineC is None: sp.line.fill.background()
    else: sp.line.color.rgb=C(lineC); sp.line.width=Pt(lw)
    return sp


def set_alpha(shape,pct):
    srgb=shape._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    srgb.append(srgb.makeelement(qn('a:alpha'),{'val':str(int(pct*1000))}))


def poly(s,pts,lineC,lw=1.5,fill=None,dash=None,fill_alpha=None):
    fb=s.shapes.build_freeform(pts[0][0],pts[0][1],scale=EMU_IN)
    fb.add_line_segments([(p[0],p[1]) for p in pts[1:]],close=True)
    sp=fb.convert_to_shape(); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
        if fill_alpha is not None: set_alpha(sp,fill_alpha)
    sp.line.color.rgb=C(lineC); sp.line.width=Pt(lw)
    if dash: sp.line.dash_style=dash
    return sp


def hbar(s,x,y,w,h,frac,fill,track=CARD2):
    rect(s,x,y,w,h,track,radius=0.5)
    if frac>0: rect(s,x,y,max(0.05,w*frac),h,fill,radius=0.5)


def eyebrow(s,t,x=0.7,y=0.5): text(s,t,x,y,8,0.3,12,BLUE,bold=True)
def title(s,t,x=0.7,y=0.82,w=12.2,size=29): text(s,t,x,y,w,0.7,size,WHITE,bold=True)


DIRS=[(0,-1),(0.866,-0.5),(0.866,0.5),(0,1),(-0.866,0.5),(-0.866,-0.5)]
def rpts(cx,cy,R,vals): return [(cx+R*v*d[0],cy+R*v*d[1]) for v,d in zip(vals,DIRS)]


def radar(s,cx,cy,R,va,vb,labels,lcols=None,fs=10):
    for lv in (0.25,0.5,0.75,1.0): poly(s,rpts(cx,cy,R,[lv]*6),BORDER,0.75)
    for gx,gy in rpts(cx,cy,R,[1.0]*6): line(s,cx,cy,gx,gy,BORDER,0.75)
    poly(s,rpts(cx,cy,R,va),BLUE,2.0,fill=BLUE,fill_alpha=28)
    poly(s,rpts(cx,cy,R,vb),GREY,1.6,dash=MSO_LINE_DASH_STYLE.DASH)
    lcols=lcols or [MUTE]*6
    for i,(d,lab) in enumerate(zip(DIRS,labels)):
        lx=cx+R*1.3*d[0]; ly=cy+R*1.3*d[1]; bw=1.6
        by=ly-0.12
        if d[1]<-0.4: by=ly-0.24
        if d[1]>0.4: by=ly+0.0
        text(s,lab,lx-bw/2,by,bw,0.24,fs,lcols[i],bold=True,align=PP_ALIGN.CENTER)


def minilag(s,x,y,scale=1.0):
    base=y+0.62*scale; vals=[0.12,0.25,0.5,1.0,0.66,0.4,0.22]
    for i,v in enumerate(vals):
        hx=x+i*0.30*scale; col=BLUE if i==3 else FAINT
        rect(s,hx,base-0.55*scale*v,0.2*scale,0.55*scale*v,col,radius=0.0,shape=MSO_SHAPE.RECTANGLE)


def minichain(s,x,y):
    pts=[x,x+0.6,x+1.2]
    line(s,pts[0]+0.13,y,pts[1]-0.13,y,BLUE,2.0)
    line(s,pts[1]+0.13,y,pts[2]-0.13,y,RED,2.0,MSO_LINE_DASH_STYLE.DASH)
    cols=[CARD2,CARD2,CARD2]
    oval(s,pts[0],y,0.13,CARD2,BORDER,1.0)
    oval(s,pts[1],y,0.15,CARD2,BORDER,1.0)
    oval(s,pts[2],y,0.13,CARD2,BORDER,1.0)
    text(s,"不在",pts[1]-0.25,y+0.18,0.5,0.2,9,FAINT,align=PP_ALIGN.CENTER)


# ===================================================== SLIDE 1: 主張
s=slide()
eyebrow(s,"提案")
title(s,"選手の「なぜ必要か」を、因果で言える")
text(s,"“効いた量（相関）”ではなく、“なぜ効くか（因果）”を1枚のレポートで示す。",0.7,1.55,12,0.4,15,MUTE)
cw,gap,x0,y0=5.9,0.4,0.7,2.35
# card 1
x=x0
rect(s,x,y0,cw,3.0,CARD,BORDER,1.0,radius=0.05)
rect(s,x,y0,0.13,3.0,BLUE,shape=MSO_SHAPE.RECTANGLE)
text(s,"① なぜ必要かを言える",x+0.4,y0+0.32,cw-0.7,0.4,16,BLUE,bold=True)
text(s,"−34%",x+0.4,y0+0.95,3,0.7,34,WHITE,bold=True)
text(s,"この選手が不在なら、チーム創出脅威",x+0.4,y0+1.72,cw-0.7,0.35,13,MUTE)
text(s,"・ 得点期待 0.34 → 0.06（連鎖が断絶）\n・ 今季の得点関与の 82% を経由（媒介中心性）",
     x+0.4,y0+2.12,cw-0.7,0.8,13.5,INK,spacing=1.3)
# card 2
x=x0+cw+gap
rect(s,x,y0,cw,3.0,CARD,BORDER,1.0,radius=0.05)
rect(s,x,y0,0.13,3.0,TEAL,shape=MSO_SHAPE.RECTANGLE)
text(s,"② 因果関係である",x+0.4,y0+0.32,cw-0.7,0.4,16,TEAL,bold=True)
text(s,"動き(t) → +0.5秒で味方の脅威↑",x+0.4,y0+1.0,cw-0.7,0.6,21,WHITE,bold=True,spacing=1.0)
text(s,"時系列で“原因 → 結果”の方向を特定",x+0.4,y0+1.72,cw-0.7,0.35,13,MUTE)
text(s,"・ 同時に起きた“相関”ではない\n・ 不在を仮想して効果を検証（反実仮想）",
     x+0.4,y0+2.12,cw-0.7,0.8,13.5,INK,spacing=1.3)
rect(s,0.7,5.75,12.2,0.95,PANEL,BORDER,1.0,radius=0.06)
text(s,"既存ツール＝相関の絵（合計 xT・VAEP・レーダー）。 Pitch Log＝因果の根拠を、納品レポート1枚に。",
     1.05,6.0,11.6,0.5,14,INK,anchor=MSO_ANCHOR.MIDDLE)

# ===================================================== SLIDE 2: レポート
s=slide()
text(s,"アウトプット：選手評価レポート（サンプル）",0.6,0.32,10,0.3,12,BLUE,bold=True)
# header band
rect(s,0.6,0.72,12.13,0.92,PANEL,BORDER,1.0,radius=0.05)
rect(s,0.6,0.72,0.13,0.92,BLUE,shape=MSO_SHAPE.RECTANGLE)
text(s,"M. Tanaka — 選手評価レポート",0.95,0.84,8,0.4,18,WHITE,bold=True)
text(s,"セントラル MF ・ 24歳 ・ FC Example ・ 推定市場価値 €12.5M",0.95,1.26,8.5,0.3,12,MUTE)
text(s,"総合 87",9.7,0.82,1.7,0.5,22,BLUE,bold=True)
rect(s,11.0,0.92,1.55,0.5,"1E5A3A",radius=0.25)
text(s,"獲得推奨",11.0,1.0,1.55,0.32,13,"9AF0C4",bold=True,align=PP_ALIGN.CENTER)
# left radar panel
rect(s,0.6,1.8,3.75,3.95,CARD,BORDER,1.0,radius=0.05)
text(s,"因果寄与プロファイル",0.85,1.98,3.3,0.3,12.5,WHITE,bold=True)
radar(s,2.45,3.72,1.02,
      [0.88,0.78,0.85,0.92,0.82,0.5],[0.6,0.62,0.58,0.5,0.5,0.55],
      ["攻撃因果","前進因果","連鎖ハブ性","オフボール牽引","代替不可性","守備因果"],
      lcols=[MUTE,MUTE,BLUE,BLUE,BLUE,MUTE],fs=9)
text(s,"■ 本選手    ┄ 同ポジ平均",0.85,5.46,3.3,0.3,10.5,MUTE,align=PP_ALIGN.CENTER)
# panel A: why necessary
ax,ay,aw,ah=4.55,1.8,8.18,1.85
rect(s,ax,ay,aw,ah,CARD,BORDER,1.0,radius=0.05)
text(s,"なぜ必要か（代替不可性）",ax+0.3,ay+0.18,5,0.3,13,WHITE,bold=True)
text(s,"−34%",ax+0.3,ay+0.62,2,0.6,28,BLUE,bold=True)
text(s,"不在時のチーム創出脅威",ax+0.32,ay+1.3,2.3,0.3,11,MUTE)
hbar(s,ax+2.45,ay+0.72,1.9,0.16,1.0,BLUE); text(s,"在籍時",ax+2.45,ay+0.52,1.5,0.2,9.5,MUTE)
hbar(s,ax+2.45,ay+1.18,1.9,0.16,0.66,GREY); text(s,"不在（仮想）",ax+2.45,ay+0.98,1.6,0.2,9.5,MUTE)
line(s,ax+4.65,ay+0.35,ax+4.65,ay+ah-0.3,BORDER,1.0)
text(s,"得点期待 0.34 → 0.06",ax+4.95,ay+0.5,3.0,0.35,14,INK,bold=True)
text(s,"経由点が抜けると連鎖が断絶",ax+4.95,ay+0.92,3.0,0.3,11,MUTE)
text(s,"媒介中心性 82%",ax+4.95,ay+1.28,3.0,0.35,14,GOLD,bold=True)
# panel B: causal basis
bx,by,bw,bh=4.55,3.8,8.18,1.95
rect(s,bx,by,bw,bh,CARD,BORDER,1.0,radius=0.05)
text(s,"因果である根拠",bx+0.3,by+0.18,5,0.3,13,WHITE,bold=True)
text(s,"時間ラグ：方向の特定",bx+0.3,by+0.6,3,0.3,11.5,TEAL,bold=True)
minilag(s,bx+0.35,by+0.92,scale=1.0)
text(s,"ピーク +0.5s ＝ 選手が先・結果が後",bx+0.3,by+1.55,3.8,0.3,10.5,MUTE)
line(s,bx+4.65,by+0.35,bx+4.65,by+bh-0.3,BORDER,1.0)
text(s,"因果クレジット配分（あるゴール）",bx+4.95,by+0.6,3.2,0.3,11.5,TEAL,bold=True)
segs=[(0.38,BLUE),(0.27,TEAL),(0.19,GREY),(0.16,FAINT)]
cx=bx+4.95; seg_w=2.9
for fr,col in segs:
    rect(s,cx,by+0.95,seg_w*fr,0.24,col,radius=0.0,shape=MSO_SHAPE.RECTANGLE); cx+=seg_w*fr
text(s,"Tanaka 38%（タッチ0のオフボール牽引）",bx+4.95,by+1.32,3.5,0.3,10.5,MUTE)
# footer conclusion
rect(s,0.6,5.9,12.13,1.0,PANEL,BORDER,1.0,radius=0.05)
rect(s,0.6,5.9,0.13,1.0,TEAL,shape=MSO_SHAPE.RECTANGLE)
text(s,"結論：契約更新／獲得を推奨",0.95,6.02,6,0.4,15,WHITE,bold=True)
text(s,"代替困難な“因果ハブ”。不在ならチームの得点創出が大きく低下するため、査定 €12.5M は妥当。",
     0.95,6.45,11.6,0.35,13,INK)

# ===================================================== SLIDE 3: 因果の根拠
s=slide()
eyebrow(s,"因果の担保")
title(s,"なぜ“因果”と言い切れるのか")
text(s,"相関では出せない3つの証拠で裏づける。",0.7,1.55,12,0.4,15,MUTE)
cards=[("時間ラグ","方向の特定","ピークが +0.5秒側＝\n選手が先・結果が後。\n同時相関ではない。",TEAL),
       ("反実仮想","介入で検証","不在を仮想して比較（−34%）。\n観察ではなく“介入効果”。",BLUE),
       ("連鎖の経路","経路の必須性","経由点を抜くと後段が\n成立しない（媒介中心性82%）。",GOLD)]
cw,gap,x0,y0=3.85,0.4,0.7,2.35
for i,(tag,concept,desc,col) in enumerate(cards):
    x=x0+i*(cw+gap)
    rect(s,x,y0,cw,3.05,CARD,BORDER,1.0,radius=0.05)
    rect(s,x,y0,cw,0.1,col,shape=MSO_SHAPE.RECTANGLE)
    text(s,tag,x+0.35,y0+0.3,cw-0.7,0.35,13,col,bold=True)
    text(s,concept,x+0.35,y0+0.66,cw-0.7,0.4,18,WHITE,bold=True)
    # mini visual
    if i==0: minilag(s,x+0.4,y0+1.25,scale=1.1)
    elif i==1:
        hbar(s,x+0.4,y0+1.35,cw-1.0,0.16,1.0,BLUE); text(s,"在籍",x+0.4,y0+1.15,1.5,0.2,9.5,MUTE)
        hbar(s,x+0.4,y0+1.75,cw-1.0,0.16,0.66,GREY); text(s,"不在",x+0.4,y0+1.55,1.5,0.2,9.5,MUTE)
    else: minichain(s,x+0.55,y0+1.55)
    text(s,desc,x+0.35,y0+2.15,cw-0.7,0.8,13,MUTE,spacing=1.25)
rect(s,0.7,5.75,12.2,0.95,PANEL,BORDER,1.0,radius=0.06)
text(s,"いずれも相関では出せない。だから“なぜ必要か”を、根拠を持って言える。",
     1.05,6.0,11.6,0.5,14,INK,anchor=MSO_ANCHOR.MIDDLE)

out=os.path.join(os.path.dirname(__file__),"Pitch_Log_Report_3p.pptx")
prs.save(out)
print("saved ->",out)
