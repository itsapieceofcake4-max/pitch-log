# -*- coding: utf-8 -*-
"""
make_report_v2.py
=================
PitchLog 提案レポート（MAZDAテンプレ様式・白/緑/游ゴシック・10x5.625in）。
出力: docs/Pitch_Log_Report_v2.pptx
構成: 1=主張 2=攻撃編レポート（不在で連鎖が途切れる画） 3=守備編レポート 4=因果の根拠
テンプレ準拠: docs/PitchLog_アンケート評価方法_テンプレ.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

# ---- テンプレ デザイントークン ----
WHITE="FFFFFF"; GREEN="1E5C3F"; GREENLT="EAF3EE"; INK="1F2A26"; MUTE="5A6663"
FAINT="9AA5A1"; GRAYROW="F6F7F6"; BORDER="D7E2DC"; RED="C0443A"; REDLT="F7ECEA"
JP="ＭＳ Ｐゴシック"


def C(h): return RGBColor.from_string(h)


prs=Presentation(); prs.slide_width=Inches(10); prs.slide_height=Inches(5.625)
BLANK=prs.slide_layouts[6]


def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=C(bg); r.line.fill.background(); r.shadow.inherit=False
    return s


def text(s,t,x,y,w,h,size,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,spacing=1.0):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,ln in enumerate(t.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=spacing
        r=p.add_run(); r.text=ln
        f=r.font; f.size=Pt(size); f.bold=bold; f.name=JP; f.color.rgb=C(color)
        rPr=r._r.get_or_add_rPr()
        for tag in ('a:ea','a:cs'):
            el=rPr.find(qn(tag))
            if el is None:
                el=rPr.makeelement(qn(tag),{}); rPr.append(el)
            el.set('typeface',JP)
    return tb


def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,radius=0.08,shape=MSO_SHAPE.ROUNDED_RECTANGLE,dash=None):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h)); sp.shadow.inherit=False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
    if line is None: sp.line.fill.background()
    else:
        sp.line.color.rgb=C(line); sp.line.width=Pt(lw)
        if dash: sp.line.dash_style=dash
    try:
        if shape==MSO_SHAPE.ROUNDED_RECTANGLE: sp.adjustments[0]=radius
    except Exception: pass
    return sp


def line(s,x1,y1,x2,y2,color=BORDER,w=1.0,dash=None):
    ln=s.shapes.add_connector(2,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    ln.line.color.rgb=C(color); ln.line.width=Pt(w); ln.shadow.inherit=False
    if dash: ln.line.dash_style=dash
    return ln


def oval(s,cx,cy,r,fill,lineC=None,lw=1.0):
    sp=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cx-r),Inches(cy-r),Inches(2*r),Inches(2*r)); sp.shadow.inherit=False
    sp.fill.solid(); sp.fill.fore_color.rgb=C(fill)
    if lineC is None: sp.line.fill.background()
    else: sp.line.color.rgb=C(lineC); sp.line.width=Pt(lw)
    return sp


def header(s,eyebrow,ttl,sub,tsize=22):
    text(s,eyebrow,0.6,0.34,8.8,0.3,11,GREEN,bold=True)
    text(s,ttl,0.6,0.6,8.8,0.55,tsize,INK,bold=True)
    text(s,sub,0.6,1.22,8.8,0.3,11,MUTE)
    text(s,"PitchLog  |  MAZDA 新規事業開発",6.6,5.32,3.2,0.25,8,FAINT,align=PP_ALIGN.RIGHT)


def hbar(s,x,y,w,h,frac,fill,track="E4EAE7"):
    rect(s,x,y,w,h,track,radius=0.5)
    if frac>0: rect(s,x,y,max(0.03,w*frac),h,fill,radius=0.5)


# nodes: list of dict {actor,label,fill,txt,lblcol,r}
# arrows: list len-1 of dict {color,dash} or None
def chain(s,yc,x0,x1,nodes,arrows):
    n=len(nodes); xs=[x0+i*(x1-x0)/(n-1) for i in range(n)]
    for i,a in enumerate(arrows):
        if a is None: continue
        r0=nodes[i].get('r',0.17); r1=nodes[i+1].get('r',0.17)
        line(s,xs[i]+r0,yc,xs[i+1]-r1,yc,a['color'],2.0,
             MSO_LINE_DASH_STYLE.DASH if a.get('dash') else None)
    for cx,nd in zip(xs,nodes):
        r=nd.get('r',0.17)
        oval(s,cx,yc,r,nd['fill'],nd.get('line'),1.0)
        text(s,nd['actor'],cx-r-0.05,yc-0.11,2*r+0.1,0.22,8,nd.get('txt',INK),
             bold=True,align=PP_ALIGN.CENTER)
        text(s,nd['label'],cx-0.7,yc+0.26,1.4,0.22,8,nd.get('lblcol',MUTE),align=PP_ALIGN.CENTER)


def conclusion(s,y,lead,body):
    rect(s,0.6,y,8.8,0.7,GREENLT,GREEN,1.0,radius=0.04)
    rect(s,0.6,y,0.06,0.7,GREEN,shape=MSO_SHAPE.RECTANGLE)
    text(s,lead,0.85,y+0.1,8.4,0.3,12,GREEN,bold=True)
    text(s,body,0.85,y+0.38,8.4,0.28,10.5,INK)


GR=lambda: dict(fill=GRAYROW,line=BORDER,txt=INK)


# ===================================================== SLIDE 1: 主張
s=slide()
header(s,"提案 ｜ 選手の貢献を“因果”で評価する",
       "選手の「なぜ必要か」を、因果で言える",
       "“効いた量（相関）”ではなく、“なぜ効くか（因果）”を1枚のレポートで示す。")
cw,gap,x0,y0=4.3,0.2,0.6,1.7
# card1
x=x0
rect(s,x,y0,cw,2.55,GREENLT,GREEN,1.0,radius=0.04)
rect(s,x,y0,0.06,2.55,GREEN,shape=MSO_SHAPE.RECTANGLE)
text(s,"① なぜ必要かを言える",x+0.3,y0+0.22,cw-0.5,0.3,13,GREEN,bold=True)
text(s,"−34%",x+0.3,y0+0.62,2,0.6,30,RED,bold=True)
text(s,"この選手が不在なら、チーム創出脅威",x+0.3,y0+1.34,cw-0.5,0.3,11,MUTE)
text(s,"・ 得点期待 0.34 → 0.06（連鎖が断絶）\n・ 得点関与の 82% を経由（媒介中心性）",
     x+0.3,y0+1.72,cw-0.5,0.7,11.5,INK,spacing=1.3)
# card2
x=x0+cw+gap
rect(s,x,y0,cw,2.55,GREENLT,GREEN,1.0,radius=0.04)
rect(s,x,y0,0.06,2.55,GREEN,shape=MSO_SHAPE.RECTANGLE)
text(s,"② 因果関係である",x+0.3,y0+0.22,cw-0.5,0.3,13,GREEN,bold=True)
text(s,"動き(t) → +0.5秒で味方の脅威↑",x+0.3,y0+0.66,cw-0.5,0.6,17,INK,bold=True,spacing=1.0)
text(s,"時系列で“原因 → 結果”の方向を特定",x+0.3,y0+1.34,cw-0.5,0.3,11,MUTE)
text(s,"・ 同時に起きた“相関”ではない\n・ 不在を仮想して効果を検証（反実仮想）",
     x+0.3,y0+1.72,cw-0.5,0.7,11.5,INK,spacing=1.3)
text(s,"既存ツール＝相関の絵（合計 xT・VAEP・レーダー）。 PitchLog＝因果の根拠を、納品レポート1枚に。",
     0.6,4.5,8.8,0.4,11,MUTE)

# ===================================================== SLIDE 2: 攻撃編
s=slide()
header(s,"アウトプット ｜ 選手評価レポート（攻撃編・サンプル）",
       "不在だと、得点への連鎖が途切れる",
       "M. Tanaka ・ セントラルMF ・ 総合 87 ・ 推定 €12.5M ・ 判定：獲得推奨")
text(s,"在籍：連鎖が通り、得点へ",0.6,1.62,5,0.25,11,GREEN,bold=True)
chain(s,2.25,1.2,8.6,
      [dict(actor="CM",label="奪回",**GR()),
       dict(actor="FB",label="配球",**GR()),
       dict(actor="田中",label="オフボール牽引",fill=GREEN,txt=WHITE,lblcol=GREEN,r=0.21),
       dict(actor="A",label="スルーパス",**GR()),
       dict(actor="FW",label="シュート",fill=GREENLT,line=GREEN,txt=GREEN,lblcol=GREEN)],
      [dict(color=GREEN)]*4)
text(s,"→ 得点期待 0.34（得点）",1.0,2.66,5,0.25,10.5,INK,bold=True)
line(s,0.6,3.02,9.4,3.02,BORDER,1.0)
text(s,"不在（反実仮想）：経由点が消え、連鎖が断絶",0.6,3.12,6,0.25,11,RED,bold=True)
chain(s,3.78,1.2,8.6,
      [dict(actor="CM",label="奪回",**GR()),
       dict(actor="FB",label="配球",**GR()),
       dict(actor="—",label="不在",fill=GRAYROW,line=BORDER,txt=FAINT,lblcol=FAINT,r=0.21),
       dict(actor="A",label="パス先消失",fill=REDLT,line=RED,txt=RED,lblcol=RED),
       dict(actor="FW",label="機会消失",fill=REDLT,line=RED,txt=RED,lblcol=RED)],
      [dict(color=MUTE),dict(color=MUTE),dict(color=RED,dash=True),dict(color=RED,dash=True)])
text(s,"→ 得点期待 0.06（−82%・断絶）",1.0,4.18,5,0.25,10.5,RED,bold=True)
conclusion(s,4.6,"結論：獲得推奨",
           "不在ならチーム創出脅威 −34% ／ 得点関与の 82% を経由する代替困難な“因果ハブ”。査定 €12.5M は妥当。")

# ===================================================== SLIDE 3: 守備編
s=slide()
header(s,"アウトプット ｜ 選手評価レポート（守備編・サンプル）",
       "不在だと、守備ブロックが空く",
       "S. Yamamoto ・ センターバック ・ 総合 84 ・ 推定 €9.0M ・ 判定：残留必須")
text(s,"在籍：カバーが効き、相手の前進を遮断",0.6,1.62,6,0.25,11,GREEN,bold=True)
chain(s,2.25,1.2,8.6,
      [dict(actor="相手",label="前進",**GR()),
       dict(actor="相手",label="楔パス",**GR()),
       dict(actor="山本",label="遮断",fill=GREEN,txt=WHITE,lblcol=GREEN,r=0.21),
       dict(actor="味方",label="回収",fill=GREENLT,line=GREEN,txt=GREEN,lblcol=GREEN),
       dict(actor="—",label="脅威なし",fill=GREENLT,line=GREEN,txt=GREEN,lblcol=GREEN)],
      [dict(color=MUTE),dict(color=MUTE),dict(color=GREEN),dict(color=GREEN)])
text(s,"→ 被xT 0.05（抑制）",1.0,2.66,5,0.25,10.5,INK,bold=True)
line(s,0.6,3.02,9.4,3.02,BORDER,1.0)
text(s,"不在：ライン間が空き、中央を割られる",0.6,3.12,6,0.25,11,RED,bold=True)
chain(s,3.78,1.2,8.6,
      [dict(actor="相手",label="前進",**GR()),
       dict(actor="相手",label="楔パス",**GR()),
       dict(actor="—",label="不在",fill=GRAYROW,line=BORDER,txt=FAINT,lblcol=FAINT,r=0.21),
       dict(actor="相手",label="中央侵入",fill=REDLT,line=RED,txt=RED,lblcol=RED),
       dict(actor="相手",label="被シュート",fill=REDLT,line=RED,txt=RED,lblcol=RED)],
      [dict(color=MUTE),dict(color=MUTE),dict(color=RED,dash=True),dict(color=RED,dash=True)])
text(s,"→ 被xT 0.31（+82%・失点機会）",1.0,4.18,5,0.25,10.5,RED,bold=True)
conclusion(s,4.6,"結論：残留必須",
           "不在なら被創出脅威 +38% ／ 危険地帯への侵入を 72% で阻止する守備の要。放出は守備崩壊リスク大。")

# ===================================================== SLIDE 4: 因果の根拠
s=slide()
header(s,"因果の担保 ｜ なぜ“相関”ではないと言えるか",
       "なぜ“因果”と言い切れるのか",
       "相関では出せない3つの証拠で裏づける。")
cards=[("時間ラグ","方向の特定","ピークが +0.5秒側＝選手が先・\n結果が後。同時相関ではない。"),
       ("反実仮想","介入で検証","不在を仮想して比較（−34%）。\n観察ではなく“介入効果”。"),
       ("連鎖の経路","経路の必須性","経由点を抜くと後段が成立\nしない（媒介中心性 82%）。")]
cw,gap,x0,y0=2.83,0.155,0.6,1.7
for i,(tag,concept,desc) in enumerate(cards):
    x=x0+i*(cw+gap)
    rect(s,x,y0,cw,2.5,GREENLT,GREEN,1.0,radius=0.04)
    rect(s,x,y0,cw,0.07,GREEN,shape=MSO_SHAPE.RECTANGLE)
    text(s,tag,x+0.28,y0+0.22,cw-0.5,0.28,11,GREEN,bold=True)
    text(s,concept,x+0.28,y0+0.54,cw-0.5,0.35,16,INK,bold=True)
    if i==0:
        base=y0+1.5; vals=[0.12,0.25,0.5,1.0,0.66,0.4,0.22]
        for j,v in enumerate(vals):
            hx=x+0.3+j*0.27; col=GREEN if j==3 else "BcC8C2".replace('Bc','B0')
            rect(s,hx,base-0.42*v,0.17,0.42*v,GREEN if j==3 else "C2CEC8",radius=0.0,shape=MSO_SHAPE.RECTANGLE)
    elif i==1:
        hbar(s,x+0.3,y0+1.15,cw-0.85,0.14,1.0,GREEN); text(s,"在籍",x+0.3,y0+0.96,1,0.2,9,MUTE)
        hbar(s,x+0.3,y0+1.5,cw-0.85,0.14,0.66,"B6C2BC"); text(s,"不在",x+0.3,y0+1.31,1,0.2,9,MUTE)
    else:
        cyc=y0+1.35; xs=[x+0.5,x+1.2,x+1.9]
        line(s,xs[0]+0.12,cyc,xs[1]-0.14,cyc,GREEN,2.0)
        line(s,xs[1]+0.14,cyc,xs[2]-0.12,cyc,RED,2.0,MSO_LINE_DASH_STYLE.DASH)
        oval(s,xs[0],cyc,0.12,GRAYROW,BORDER,1.0)
        oval(s,xs[1],cyc,0.14,GRAYROW,BORDER,1.0)
        oval(s,xs[2],cyc,0.12,REDLT,RED,1.0)
        text(s,"不在",xs[1]-0.25,cyc+0.16,0.5,0.2,8,FAINT,align=PP_ALIGN.CENTER)
    text(s,desc,x+0.28,y0+1.92,cw-0.5,0.55,10.5,MUTE,spacing=1.2)
text(s,"いずれも相関では出せない。だから“なぜ必要か”を、根拠を持って言える。",
     0.6,4.45,8.8,0.4,11,MUTE)

out=os.path.join(os.path.dirname(__file__),"Pitch_Log_Report_v2.pptx")
prs.save(out)
print("saved ->",out)
