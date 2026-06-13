# -*- coding: utf-8 -*-
r"""
make_realdata_report.py
=======================
xT付与済みトラッキングCSV(soccernet_add_xt.py の出力)から、
選手別の貢献度(xT)を集計し、MAZDAテンプレ様式の実データ図(pptx)を出す。

  python make_realdata_report.py --csv <...>_tracking_xt.csv --team home
  → 集計を表示し、docs/Pitch_Log_RealData_<clip>.pptx を生成。
"""
import argparse
import os
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

WHITE="FFFFFF"; GREEN="1E5C3F"; GREENLT="EAF3EE"; INK="1F2A26"; MUTE="5A6663"
FAINT="9AA5A1"; BORDER="D7E2DC"; GRAYBAR="C2CEC8"; JP="ＭＳ Ｐゴシック"


def C(h): return RGBColor.from_string(h)


def _font(r):
    r.font.name = JP
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", JP)


def text(s, t, x, y, w, h, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = t
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = C(color); _font(r)
    return tb


def rect(s, x, y, w, h, fill=None, ln=None, lw=1.0, radius=0.06,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if ln is None: sp.line.fill.background()
    else: sp.line.color.rgb = C(ln); sp.line.width = Pt(lw)
    try:
        if shape == MSO_SHAPE.ROUNDED_RECTANGLE: sp.adjustments[0] = radius
    except Exception: pass
    return sp


def conn(s, x1, y1, x2, y2, color, w=1.5):
    c = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = C(color); c.line.width = Pt(w); c.shadow.inherit = False
    return c


def aggregate(df, clip, team):
    d = df[(df["clip"] == clip) & (df["team"] == team) & df["xt"].notna()]
    g = d.groupby("track_id")
    rows = []
    for tid, sub in g:
        js = sub["jersey"].dropna()
        jersey = "" if js.empty else str(int(js.mode().iloc[0])) if str(js.mode().iloc[0]).replace('.0','').isdigit() else str(js.mode().iloc[0])
        rows.append({"track_id": tid, "jersey": jersey,
                     "frames": len(sub), "mean_xt": sub["xt"].mean(),
                     "peak_xt": sub["xt"].max()})
    r = pd.DataFrame(rows).sort_values("mean_xt", ascending=False).reset_index(drop=True)
    # チーム脅威タイムライン（フレームごとの xt 合計）
    cd = df[df["clip"] == clip]
    fmin, fmax = int(cd["frame"].min()), int(cd["frame"].max())
    tl = d.groupby("frame")["xt"].sum().reindex(range(fmin, fmax + 1), fill_value=0.0)
    return r, tl


def build(clip, action, rank, tl, team, out):
    prs = Presentation(); prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = C(WHITE); bg.line.fill.background(); bg.shadow.inherit = False

    text(s, f"実データ ｜ SoccerNet GSR — {clip}", 0.6, 0.34, 8.8, 0.3, 11, GREEN, bold=True)
    text(s, "選手別 貢献度（xT）— 実データ検証", 0.6, 0.6, 8.8, 0.5, 20, INK, bold=True)
    text(s, f"アクション：{action} ・ 30秒/750frame ・ 攻撃方向はGKから自動推定（{team}）",
         0.6, 1.18, 8.8, 0.3, 11, MUTE)
    text(s, "PitchLog  |  MAZDA 新規事業開発", 6.6, 5.32, 3.2, 0.25, 8, FAINT, align=PP_ALIGN.RIGHT)

    # 左：ランキング
    rect(s, 0.6, 1.7, 5.45, 3.55, GREENLT, BORDER, 1.0, radius=0.04)
    text(s, "平均xT 上位（攻撃側）", 0.85, 1.88, 5, 0.3, 12.5, GREEN, bold=True)
    top = rank.head(8).reset_index(drop=True)
    vmax = max(top["mean_xt"].max(), 1e-6)
    bx, bw, y = 2.05, 3.0, 2.35
    for _, r in top.iterrows():
        lab = f"背番号 {r['jersey']}" if r["jersey"] else f"ID {int(r['track_id'])}"
        text(s, lab, 0.85, y - 0.02, 1.2, 0.25, 10, INK)
        rect(s, bx, y + 0.02, bw, 0.16, "DDE7E1", radius=0.5)
        rect(s, bx, y + 0.02, max(0.04, bw * r["mean_xt"] / vmax), 0.16, GREEN, radius=0.5)
        text(s, f"{r['mean_xt']:.2f}", bx + bw + 0.08, y - 0.02, 0.6, 0.25, 10, INK, bold=True)
        y += 0.345

    # 右：クリップ脅威タイムライン
    rect(s, 6.2, 1.7, 3.2, 3.55, WHITE, BORDER, 1.0, radius=0.04)
    text(s, "チーム脅威の推移（xT合計）", 6.42, 1.88, 3, 0.3, 12.5, GREEN, bold=True)
    vals = tl.to_numpy(dtype=float)
    n = len(vals); step = max(1, n // 48)
    ds = vals[::step]
    x0, x1, ytop, ybot = 6.5, 9.1, 2.45, 4.55
    vmx = max(ds.max(), 1e-6)
    pts = [(x0 + (x1 - x0) * i / (len(ds) - 1), ybot - (ybot - ytop) * (v / vmx))
           for i, v in enumerate(ds)]
    conn(s, x0, ybot, x1, ybot, BORDER, 1.0)
    for i in range(len(pts) - 1):
        conn(s, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], GREEN, 1.75)
    pk = int(np.argmax(ds))
    text(s, "ピーク＝コーナー供給時", 6.5, ytop - 0.16, 3, 0.25, 9.5, MUTE)
    text(s, "0秒", 6.4, 4.62, 1, 0.2, 9, FAINT)
    text(s, "30秒", 8.7, 4.62, 1, 0.2, 9, FAINT, align=PP_ALIGN.RIGHT)
    text(s, f"クリップ最大 xT合計：{ds.max():.2f}", 6.42, 4.9, 3, 0.25, 10, INK, bold=True)

    text(s, "※ xT＝位置脅威（攻撃方向で向き付け）。放送由来トラッキングのため未検出フレームあり。"
         "本図は実データ(SoccerNet GSR)での貢献度可視化のサンプル。",
         0.6, 5.32, 6.0, 0.3, 8.5, FAINT)

    prs.save(out)
    return out


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--csv", required=True)
    ap.add_argument("--team", default="home")
    ap.add_argument("--clip", default=None)
    ap.add_argument("--action", default="コーナー")
    ap.add_argument("--out", default=None)
    args = ap.parse_args()

    df = pd.read_csv(args.csv)
    df["team"] = df["team"].fillna("")
    clip = args.clip or df["clip"].iloc[0]
    rank, tl = aggregate(df, clip, args.team)
    print(f"=== {clip} / {args.team} 貢献度ランキング（平均xT） ===")
    print(rank.head(8).to_string(index=False))

    out = args.out or os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                   f"Pitch_Log_RealData_{clip}.pptx")
    build(clip, args.action, rank, tl, args.team, out)
    print("saved ->", out)


if __name__ == "__main__":
    main()
