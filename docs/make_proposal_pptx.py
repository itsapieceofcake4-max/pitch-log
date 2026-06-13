# -*- coding: utf-8 -*-
"""
make_proposal_pptx.py
=====================
Pitch Log 顧客提案スライドを生成する。
出力: docs/Pitch_Log_Proposal.pptx  (16:9)
構成: ①レーダー → 因果6軸 → ②因果ビュー → 連鎖断絶 → ③シミュレータ → 差別化 → まとめ
"""
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

# ---- ブランドカラー ----
BG     = "0E1825"
PANEL  = "16273B"
CARD   = "1A2B41"
CARD2  = "203busy"  # placeholder, replaced below
CARD2  = "21364F"
BORDER = "2C4A6E"
WHITE  = "FFFFFF"
INK    = "E8F0FA"
MUTE   = "93AAC6"
FAINT  = "647E9B"
BLUE   = "3B9EFF"
BLUEDK = "1C3F66"
TEAL   = "2BD4A8"
RED    = "FF5C5C"
REDDK  = "63242A"
GOLD   = "FFC24B"
GREY   = "8FA0B4"
JP = "Meiryo"

EMU_IN = 914400


def C(h):
    return RGBColor.from_string(h)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = C(bg)
    r.line.fill.background(); r.shadow.inherit = False
    return s


def text(s, t, x, y, w, h, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font=JP, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = t.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.name = font
        f.color.rgb = C(color)
    return tb


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, radius=0.08, shape=MSO_SHAPE.ROUNDED_RECTANGLE, dash=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(line); sp.line.width = Pt(lw)
        if dash:
            sp.line.dash_style = dash
    try:
        if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
            sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def line(s, x1, y1, x2, y2, color=BORDER, w=1.0, dash=None):
    ln = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = C(color); ln.line.width = Pt(w)
    ln.shadow.inherit = False
    if dash:
        ln.line.dash_style = dash
    return ln


def oval(s, cx, cy, r, fill, lineC=None, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r))
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
    if lineC is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = C(lineC); sp.line.width = Pt(lw)
    return sp


def set_fill_alpha(shape, pct):
    srgb = shape._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    al = srgb.makeelement(qn('a:alpha'), {'val': str(int(pct * 1000))})
    srgb.append(al)


def poly(s, pts, lineC, lw=1.5, fill=None, dash=None, fill_alpha=None):
    fb = s.shapes.build_freeform(pts[0][0], pts[0][1], scale=EMU_IN)
    fb.add_line_segments([(p[0], p[1]) for p in pts[1:]], close=True)
    sp = fb.convert_to_shape()
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = C(fill)
        if fill_alpha is not None:
            set_fill_alpha(sp, fill_alpha)
    sp.line.color.rgb = C(lineC); sp.line.width = Pt(lw)
    if dash:
        sp.line.dash_style = dash
    return sp


def eyebrow(s, t, x=0.7, y=0.55):
    text(s, t, x, y, 6, 0.3, 12, BLUE, bold=True)


def title(s, t, x=0.7, y=0.85, w=12, size=27):
    text(s, t, x, y, w, 0.7, size, WHITE, bold=True)


DIRS = [(0, -1), (0.866, -0.5), (0.866, 0.5), (0, 1), (-0.866, 0.5), (-0.866, -0.5)]


def radar_pts(cx, cy, R, vals):
    return [(cx + R * v * d[0], cy + R * v * d[1]) for v, d in zip(vals, DIRS)]


def radar(s, cx, cy, R, vals_a, vals_b, labels, label_colors=None):
    for lv in (0.25, 0.5, 0.75, 1.0):
        gp = radar_pts(cx, cy, R, [lv] * 6)
        poly(s, gp, BORDER, 0.75)
    grid = radar_pts(cx, cy, R, [1.0] * 6)
    for gx, gy in grid:
        line(s, cx, cy, gx, gy, BORDER, 0.75)
    pa = radar_pts(cx, cy, R, vals_a)
    poly(s, pa, BLUE, 2.25, fill=BLUE, fill_alpha=28)
    pb = radar_pts(cx, cy, R, vals_b)
    poly(s, pb, GREY, 1.75, dash=MSO_LINE_DASH_STYLE.DASH)
    label_colors = label_colors or [MUTE] * 6
    for i, (d, lab) in enumerate(zip(DIRS, labels)):
        lx = cx + R * 1.28 * d[0]
        ly = cy + R * 1.28 * d[1]
        align = PP_ALIGN.CENTER
        bw = 1.7
        bx = lx - bw / 2
        by = ly - 0.13
        if d[1] < -0.4:
            by = ly - 0.26
        if d[1] > 0.4:
            by = ly + 0.02
        text(s, lab, bx, by, bw, 0.26, 11, label_colors[i], bold=True, align=align)


def hbar(s, x, y, w, h, frac, fill, track=CARD2, radius=0.5):
    rect(s, x, y, w, h, track, radius=radius)
    if frac > 0:
        rect(s, x, y, max(0.05, w * frac), h, fill, radius=radius)


def legend(s, x, y, items):
    cx = x
    for lab, col, dash in items:
        if dash:
            line(s, cx, y + 0.1, cx + 0.28, y + 0.1, col, 2.0, MSO_LINE_DASH_STYLE.DASH)
        else:
            rect(s, cx, y + 0.03, 0.28, 0.13, col, radius=0.3)
        text(s, lab, cx + 0.36, y - 0.04, 1.6, 0.26, 11, MUTE)
        cx += 0.42 + 0.13 * len(lab) + 0.4


# ============================================================ Slide 1: Title
s = slide()
rect(s, 0, 0, 0.22, 7.5, BLUE, shape=MSO_SHAPE.RECTANGLE)
text(s, "PITCH LOG", 0.9, 2.0, 6, 0.4, 15, BLUE, bold=True)
text(s, "因果で測る、選手の本当の価値", 0.9, 2.5, 11.5, 1.2, 40, WHITE, bold=True)
text(s, "戦術分析 → 貢献度 → 査定 → 編成・移籍 をつなぐ意思決定基盤", 0.95, 3.85, 11, 0.5, 18, MUTE)
line(s, 0.95, 4.55, 5.0, 4.55, BORDER, 1.0)
text(s, "サッカー戦術データ分析プラットフォーム ｜ 顧客提案資料", 0.95, 4.7, 11, 0.4, 13, FAINT)
text(s, "© Pitch Log", 0.95, 6.7, 6, 0.3, 11, FAINT)

# ============================================================ Slide 2: Problem
s = slide()
eyebrow(s, "課題")
title(s, "既存の評価は「相関」で止まっている")
probs = [
    ("ti", "記述指標の限界", "合計 xT・VAEP・走行距離…\n“やった量”は分かる。"),
    ("ti", "「なぜ効くか」が不明", "レーダーやランキングは相関。\n因果は語れない。"),
    ("ti", "「抜けたら」が不明", "選手の必要性・代替可能性を\n定量化できない。"),
]
cw, gap, x0, y0 = 3.7, 0.45, 0.9, 2.1
for i, (_, h, d) in enumerate(probs):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, 2.7, CARD, BORDER, 1.0, radius=0.06)
    rect(s, x + 0.35, y0 + 0.4, 0.5, 0.08, RED, radius=0.3)
    text(s, h, x + 0.35, y0 + 0.65, cw - 0.7, 0.6, 18, WHITE, bold=True)
    text(s, d, x + 0.35, y0 + 1.45, cw - 0.7, 1.1, 14, MUTE, spacing=1.15)
rect(s, 0.9, 5.25, 11.5, 0.9, PANEL, BORDER, 1.0, radius=0.08)
text(s, "業界標準＝StatsBomb レーダー＋動画。Wyscout・SkillCorner も含め、ほぼ全社が同じ“相関の絵”。",
     1.2, 5.5, 11, 0.5, 14, INK, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ Slide 3: Thesis
s = slide()
eyebrow(s, "提案の核")
title(s, "相関から、因果へ")
text(s, "“効いた量”ではなく、“引き起こした価値”を測る。", 0.9, 1.6, 11, 0.4, 15, MUTE)
pills = [
    ("反実仮想", "居なければチーム脅威 −34%", "不在を仮想して価値を測る", BLUE),
    ("時間ラグ因果", "動き(t) → 0.5秒後の味方脅威", "“引き起こした”方向を特定", TEAL),
    ("オフボール因果クレジット", "タッチ0でも貢献を配分", "走りでDFを引きつけた価値", GOLD),
]
cw, gap, x0, y0 = 3.7, 0.45, 0.9, 2.3
for i, (h, big, d, col) in enumerate(pills):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, 3.0, CARD, BORDER, 1.0, radius=0.06)
    text(s, "0" + str(i + 1), x + 0.35, y0 + 0.3, 1, 0.5, 16, col, bold=True)
    text(s, h, x + 0.35, y0 + 0.75, cw - 0.7, 0.5, 16, WHITE, bold=True)
    text(s, big, x + 0.35, y0 + 1.45, cw - 0.7, 0.8, 19, col, bold=True, spacing=1.05)
    text(s, d, x + 0.35, y0 + 2.4, cw - 0.7, 0.5, 13, MUTE)

# ============================================================ Slide 4: Overview
s = slide()
eyebrow(s, "アウトプットの全体像")
title(s, "個の評価から、編成の意思決定まで一気通貫")
steps = ["戦術分析", "貢献度の可視化", "選手査定", "編成・移籍"]
sub = ["局面に分解", "誰が価値を生んだか", "総合スコア化・比較", "更新／獲得／放出"]
bw, bgap, x0, y = 2.7, 0.55, 0.9, 2.2
for i, (st, sb) in enumerate(zip(steps, sub)):
    x = x0 + i * (bw + bgap)
    col = BLUE if i < 3 else TEAL
    rect(s, x, y, bw, 1.3, CARD, col, 1.25, radius=0.1)
    text(s, st, x, y + 0.28, bw, 0.5, 18, WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(s, sb, x, y + 0.82, bw, 0.4, 12, MUTE, align=PP_ALIGN.CENTER)
    if i < 3:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw + 0.05), Inches(y + 0.5), Inches(0.45), Inches(0.3))
        a.fill.solid(); a.fill.fore_color.rgb = C(FAINT); a.line.fill.background(); a.shadow.inherit = False
text(s, "これを支える3つの画面", 0.9, 4.0, 11, 0.4, 14, BLUE, bold=True)
screens = [("① 選手評価ダッシュボード", "レーダーで掴む（入口）"),
           ("② 因果ビュー", "独自の深掘り（差別化）"),
           ("③ 反実仮想シミュレータ", "編成・移籍の判断")]
cw, gap, x0, y0 = 3.7, 0.45, 0.9, 4.5
for i, (h, d) in enumerate(screens):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, 1.7, PANEL, BORDER, 1.0, radius=0.06)
    text(s, h, x + 0.3, y0 + 0.35, cw - 0.6, 0.5, 16, WHITE, bold=True)
    text(s, d, x + 0.3, y0 + 1.0, cw - 0.6, 0.5, 13, MUTE)

# ============================================================ Slide 5: Radar (descriptive)
s = slide()
eyebrow(s, "① 選手評価ダッシュボード")
title(s, "まず、既存の言語で手に取らせる")
radar(s, 3.5, 4.0, 1.4,
      [0.92, 0.78, 0.55, 0.85, 0.72, 0.6],
      [0.6, 0.62, 0.68, 0.55, 0.5, 0.7],
      ["攻撃貢献", "前進性", "守備貢献", "オフボール", "個の打開", "走力"])
legend(s, 2.1, 6.5, [("本選手", BLUE, False), ("同ポジ平均", GREY, True)])
# right: metric cards
cards = [("総合評価", "87", "/100 ・上位6%"), ("攻撃貢献 xT/90", "+0.42", ""),
         ("守備貢献 VAEP", "78", ""), ("推定市場価値", "€12.5M", "")]
x0, y0, cw, ch, gx, gy = 6.7, 2.2, 2.7, 1.55, 0.4, 0.4
for i, (lab, val, sub2) in enumerate(cards):
    x = x0 + (i % 2) * (cw + gx); y = y0 + (i // 2) * (ch + gy)
    rect(s, x, y, cw, ch, CARD, BORDER, 1.0, radius=0.08)
    text(s, lab, x + 0.3, y + 0.25, cw - 0.6, 0.4, 13, MUTE)
    text(s, val, x + 0.3, y + 0.62, cw - 0.6, 0.6, 28, BLUE, bold=True)
    if sub2:
        text(s, sub2, x + 0.3, y + 1.18, cw - 0.6, 0.3, 11, FAINT)
text(s, "ポジション別パーセンタイル・レーダー＝既存ツールと同じ“入口”。ここで掴んで、深掘りへ誘導する。",
     6.7, 5.95, 6.0, 0.8, 13, MUTE, spacing=1.2)

# ============================================================ Slide 6: Causal radar
s = slide()
eyebrow(s, "差別化の起点")
title(s, "レーダーを「因果寄与」で塗り替える")
radar(s, 3.5, 4.0, 1.4,
      [0.88, 0.78, 0.85, 0.92, 0.82, 0.5],
      [0.6, 0.62, 0.58, 0.5, 0.5, 0.55],
      ["攻撃因果", "前進因果", "連鎖ハブ性", "オフボール牽引", "代替不可性", "守備因果"],
      label_colors=[MUTE, MUTE, BLUE, BLUE, BLUE, MUTE])
legend(s, 2.1, 6.5, [("本選手", BLUE, False), ("同ポジ平均", GREY, True)])
text(s, "記述軸 → 因果軸への置き換え", 6.7, 2.2, 6, 0.4, 15, WHITE, bold=True)
maps = [("オフボール牽引", "タッチ0でDFを引く因果クレジット"),
        ("代替不可性", "不在時にチーム脅威が落ちる度合い"),
        ("連鎖ハブ性", "因果ネットワークで誰を動かすか")]
y = 2.9
for h, d in maps:
    rect(s, 6.7, y, 0.12, 0.55, BLUE, shape=MSO_SHAPE.RECTANGLE)
    text(s, h, 6.95, y, 5.6, 0.35, 15, BLUE, bold=True)
    text(s, d, 6.95, y + 0.36, 5.6, 0.35, 12.5, MUTE)
    y += 0.85
text(s, "“個の打開・走力”のような単独スキル軸を捨て、因果でしか測れない軸に入れ替える。",
     6.7, 5.6, 6.0, 0.8, 13, MUTE, spacing=1.2)

# ============================================================ Slide 7: Causal view (4 panels)
s = slide()
eyebrow(s, "② 因果ビュー")
title(s, "“なぜ効くか”を4つの角度で示す")
pw, ph, gx, gy, x0, y0 = 5.7, 1.95, 0.4, 0.35, 0.9, 2.0


def panel(x, y, t):
    rect(s, x, y, pw, ph, CARD, BORDER, 1.0, radius=0.06)
    text(s, t, x + 0.3, y + 0.2, pw - 0.6, 0.35, 13.5, WHITE, bold=True)


# p1 反実仮想
px, py = x0, y0
panel(px, py, "反実仮想：不在ならどうなるか")
text(s, "−34%", px + 0.3, py + 0.65, 2, 0.7, 30, BLUE, bold=True)
text(s, "チーム創出脅威 / 30分", px + 0.35, py + 1.4, 2.6, 0.3, 11, MUTE)
hbar(s, px + 2.7, py + 0.75, 2.6, 0.18, 1.0, BLUE)
text(s, "在籍時", px + 2.7, py + 0.55, 2, 0.2, 10, MUTE)
hbar(s, px + 2.7, py + 1.25, 2.6, 0.18, 0.66, GREY)
text(s, "不在（仮想）", px + 2.7, py + 1.05, 2, 0.2, 10, MUTE)

# p2 credit
px = x0 + pw + gx
panel(px, py, "因果クレジット配分（あるゴール）")
segs = [("Tanaka(オフボール)", 0.38, BLUE), ("配球者A", 0.27, TEAL),
        ("キャリーB", 0.19, GREY), ("起点C", 0.16, FAINT)]
bx = px + 0.3; bw2 = pw - 0.6
cx = bx
for lab, fr, col in segs:
    rect(s, cx, py + 0.7, bw2 * fr, 0.28, col, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
    cx += bw2 * fr
text(s, "Tanaka 38% ・ A 27% ・ B 19% ・ C 16%", px + 0.3, py + 1.1, pw - 0.6, 0.3, 11.5, INK)
text(s, "↳ Tanaka はこのシーンでタッチ0。走りに38%を配分。", px + 0.3, py + 1.5, pw - 0.6, 0.3, 11, BLUE)

# p3 network
px, py = x0, y0 + ph + gy
panel(px, py, "因果インフルエンス・ネットワーク")
ncx, ncy = px + 1.45, py + 1.05
nodes = [(px + 2.85, py + 0.65, "FW", BLUE, 2.5), (px + 2.9, py + 1.5, "WG", BLUE, 1.5),
         (px + 0.5, py + 0.65, "FB", GREY, 1.25), (px + 0.5, py + 1.5, "CM", GREY, 1.0)]
for nx, ny, lab, col, lw in nodes:
    if col == GREY:
        line(s, nx, ny, ncx, ncy, GREY, lw)
    else:
        line(s, ncx, ncy, nx, ny, BLUE, lw)
oval(s, ncx, ncy, 0.3, BLUE); text(s, "MF", ncx - 0.3, ncy - 0.15, 0.6, 0.3, 11, WHITE, bold=True, align=PP_ALIGN.CENTER)
for nx, ny, lab, col, lw in nodes:
    oval(s, nx, ny, 0.23, CARD2, BORDER, 1.0)
    text(s, lab, nx - 0.23, ny - 0.13, 0.46, 0.28, 10, INK, align=PP_ALIGN.CENTER)
text(s, "太さ＝因果の強さ", px + pw - 1.75, py + 0.28, 1.5, 0.25, 10, MUTE, align=PP_ALIGN.RIGHT)

# p4 time-lag
px = x0 + pw + gx
panel(px, py, "時間ラグ：因果の“方向”")
bx = px + 0.45; base = py + 1.18
lags = [0.12, 0.25, 0.5, 1.0, 0.66, 0.4, 0.22]
labs = ["-1.0", "-0.5", "0", "+0.5", "+1.0", "+1.5", "+2.0"]
for i, (v, lb) in enumerate(zip(lags, labs)):
    hx = bx + i * 0.62
    col = BLUE if i == 3 else FAINT
    hgt = 0.52 * v
    rect(s, hx, base - hgt, 0.42, hgt, col, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
    text(s, lb, hx - 0.07, base + 0.02, 0.56, 0.2, 8.5, MUTE, align=PP_ALIGN.CENTER)
text(s, "ピーク +0.5s ＝ 選手が先・味方の脅威が後", px + 0.3, py + 1.58, pw - 0.6, 0.3, 11, BLUE)

# ============================================================ Slide 8: Chain break
s = slide()
eyebrow(s, "「なぜ必要か」")
title(s, "経由点が抜けると、連鎖が断たれる")
chain = ["奪回", "配球", "オフボール牽引", "スルーパス", "シュート"]
actors = ["CM", "FB", "Tanaka", "A", "FW"]
XS = [1.4, 3.0, 4.6, 6.2, 7.8]


def chain_row(yc, removed):
    for i in range(4):
        if removed and i >= 2:
            line(s, XS[i] + 0.32, yc, XS[i + 1] - 0.32, yc, RED, 2.0, MSO_LINE_DASH_STYLE.DASH)
        else:
            line(s, XS[i] + 0.32, yc, XS[i + 1] - 0.32, yc, BLUE, 2.25)
    for i, (cxp, act, lab) in enumerate(zip(XS, actors, chain)):
        is_t = (i == 2)
        if is_t and removed:
            oval(s, cxp, yc, 0.3, CARD2, BORDER, 1.0)
            text(s, "不在", cxp - 0.3, yc - 0.13, 0.6, 0.28, 10, FAINT, bold=True, align=PP_ALIGN.CENTER)
        elif is_t:
            oval(s, cxp, yc, 0.34, BLUE)
            text(s, "Tanaka", cxp - 0.4, yc - 0.13, 0.8, 0.28, 9.5, WHITE, bold=True, align=PP_ALIGN.CENTER)
        else:
            fillc = TEAL if (i == 4 and not removed) else CARD2
            oval(s, cxp, yc, 0.28, fillc, BORDER, 1.0)
            tcol = "04342C" if (i == 4 and not removed) else INK
            text(s, act, cxp - 0.28, yc - 0.12, 0.56, 0.26, 10, tcol, align=PP_ALIGN.CENTER)
        lc = MUTE
        if removed and i >= 2:
            lab = ["", "", "—", "パス先消失", "機会消失"][i]; lc = RED
        text(s, lab, cxp - 0.7, yc + 0.44, 1.4, 0.3, 10.5, lc, align=PP_ALIGN.CENTER)


text(s, "連鎖あり", 0.9, 2.0, 2, 0.3, 13, BLUE, bold=True)
chain_row(2.75, False)
text(s, "ゴール期待 0.34 → 得点", 0.9, 3.55, 5, 0.3, 12.5, INK)
line(s, 0.9, 4.0, 8.5, 4.0, BORDER, 1.0)
text(s, "反実仮想：Tanaka を外す", 0.9, 4.2, 5, 0.3, 13, RED, bold=True)
chain_row(4.95, True)
text(s, "ゴール期待 0.06（−82%）→ 断絶", 0.9, 5.75, 6, 0.3, 12.5, RED)
rect(s, 9.4, 3.0, 3.0, 1.9, PANEL, BORDER, 1.0, radius=0.08)
text(s, "82%", 9.4, 3.25, 3.0, 0.7, 32, GOLD, bold=True, align=PP_ALIGN.CENTER)
text(s, "今季の得点関与で\nこの選手を経由（媒介中心性）", 9.6, 4.1, 2.6, 0.6, 11, MUTE, align=PP_ALIGN.CENTER, spacing=1.1)
text(s, "時系列で因果を追うからこそ、「経由点が抜けると後段が成立しない」を示せる。",
     0.9, 6.4, 11.5, 0.4, 13, MUTE)

# ============================================================ Slide 9: Simulator
s = slide()
eyebrow(s, "③ 反実仮想シミュレータ")
title(s, "“もし入れ替えたら”で編成・移籍を判断")
# pitch
rect(s, 0.9, 2.1, 4.6, 3.4, PANEL, BORDER, 1.0, radius=0.04, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
line(s, 3.2, 2.1, 3.2, 5.5, BORDER, 1.0)
oval(s, 3.2, 3.8, 0.45, PANEL, BORDER, 1.0)
home = [(1.35, 3.8), (1.95, 2.7), (1.95, 3.45), (1.95, 4.15), (1.95, 4.9),
        (3.0, 3.55), (3.0, 4.35), (4.2, 2.9), (4.2, 3.8), (4.2, 4.7)]
for hx, hy in home:
    oval(s, hx, hy, 0.13, GREY)
oval(s, 3.0, 2.75, 0.2, BLUE)
text(s, "MT", 2.8, 2.63, 0.4, 0.24, 9, WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, "CM スロットに候補を充当", 0.9, 5.55, 4.6, 0.3, 11, MUTE, align=PP_ALIGN.CENTER)
# right: result
text(s, "反実仮想の結果（チーム全体）", 6.0, 2.1, 6, 0.4, 14, WHITE, bold=True)
text(s, "1.78", 6.0, 2.55, 2, 0.7, 34, BLUE, bold=True)
text(s, "チーム期待 xT / 30分", 6.0, 3.3, 3, 0.3, 12, MUTE)
rect(s, 9.6, 2.65, 2.8, 0.5, "1E5A3A", radius=0.2)
text(s, "+25% ・ 獲得推奨", 9.6, 2.74, 2.8, 0.3, 14, "9AF0C4", bold=True, align=PP_ALIGN.CENTER)
# candidate compare
cands = [("M. Tanaka", 0.89, "87", BLUE, True), ("候補 X", 0.78, "79", GREY, False), ("候補 Y", 0.72, "72", GREY, False)]
y = 3.95
for nm, fr, sc, col, rec in cands:
    text(s, nm, 6.0, y, 1.8, 0.3, 13, INK)
    hbar(s, 7.7, y + 0.03, 3.0, 0.2, fr, col)
    text(s, sc, 10.85, y - 0.02, 0.6, 0.3, 13, WHITE, bold=True)
    if rec:
        text(s, "推奨", 11.6, y - 0.02, 0.8, 0.3, 12, TEAL, bold=True)
    y += 0.62
text(s, "②の因果モデルを“入れ替え”に適用 → チーム全体の創出脅威を再計算。フィットとコストで判断。",
     6.0, 6.0, 6.4, 0.6, 13, MUTE, spacing=1.2)

# ============================================================ Slide 10: Differentiation
s = slide()
eyebrow(s, "差別化")
title(s, "既存ツールとの違い")
rows = [
    ("観点", "既存（StatsBomb / Wyscout / SkillCorner）", "Pitch Log"),
    ("評価の根拠", "相関・記述（合計 xT / VAEP）", "因果（反実仮想・時間ラグ）"),
    ("なぜ効くか", "示せない", "因果クレジットで分解"),
    ("抜けたらどうなるか", "示せない", "連鎖断絶＋媒介中心性で定量化"),
    ("オフボールの価値", "ほぼ未評価", "タッチ0でも貢献を配分"),
    ("意思決定支援", "ランキング・類似検索", "反実仮想シミュレータで編成判断"),
]
x0, y0 = 0.9, 1.9
c0, c1, c2 = 3.0, 5.4, 3.5
rh = 0.78
for i, (a, b, c) in enumerate(rows):
    y = y0 + i * rh
    head = (i == 0)
    if head:
        rect(s, x0, y, c0, rh, PANEL, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
        rect(s, x0 + c0, y, c1, rh, PANEL, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
        rect(s, x0 + c0 + c1, y, c2, rh, "16364E", radius=0.0, shape=MSO_SHAPE.RECTANGLE)
    else:
        rect(s, x0 + c0 + c1, y, c2, rh, "11243A", radius=0.0, shape=MSO_SHAPE.RECTANGLE)
    ac = WHITE if head else MUTE
    bc = INK if head else MUTE
    cc = BLUE if head else INK
    text(s, a, x0 + 0.2, y, c0 - 0.3, rh, 13, ac, bold=head, anchor=MSO_ANCHOR.MIDDLE)
    text(s, b, x0 + c0 + 0.2, y, c1 - 0.3, rh, 12.5, bc, bold=head, anchor=MSO_ANCHOR.MIDDLE)
    text(s, c, x0 + c0 + c1 + 0.2, y, c2 - 0.3, rh, 12.5, cc, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    line(s, x0, y + rh, x0 + c0 + c1 + c2, y + rh, BORDER, 0.75)

# ============================================================ Slide 11: Closing
s = slide()
rect(s, 0, 0, 0.22, 7.5, BLUE, shape=MSO_SHAPE.RECTANGLE)
text(s, "まとめ", 0.9, 1.3, 6, 0.4, 14, BLUE, bold=True)
text(s, "因果が、査定と編成を変える", 0.9, 1.8, 11.5, 1.0, 34, WHITE, bold=True)
eff = [("査定の根拠が明確", "“居なければ−34%”で年俸・契約を説明"),
       ("編成を最適化", "因果フィットで穴とハマりを可視化"),
       ("移籍リスクを低減", "獲得前に反実仮想で効果を検証")]
cw, gap, x0, y0 = 3.7, 0.45, 0.9, 3.1
for i, (h, d) in enumerate(eff):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, 1.9, CARD, BORDER, 1.0, radius=0.06)
    text(s, h, x + 0.3, y0 + 0.35, cw - 0.6, 0.5, 16, WHITE, bold=True)
    text(s, d, x + 0.3, y0 + 1.0, cw - 0.6, 0.7, 13, MUTE, spacing=1.15)
rect(s, 0.9, 5.5, 11.5, 0.9, PANEL, BORDER, 1.0, radius=0.08)
text(s, "実証データ：PFF FC（2022 W杯・64試合）／ SoccerNet（映像付き）など公開データで PoC 可能。",
     1.2, 5.75, 11, 0.4, 13.5, INK, anchor=MSO_ANCHOR.MIDDLE)

import os
out = os.path.join(os.path.dirname(__file__), "Pitch_Log_Proposal.pptx")
prs.save(out)
print("saved ->", out)
