"""
make_catalog_pptx.py
====================
docs/feature_catalog.md の内容を PowerPoint 化する。
出力: docs/feature_catalog.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG     = RGBColor(0x0E, 0x18, 0x25)
SURF   = RGBColor(0x1A, 0x2E, 0x42)
HEADBG = RGBColor(0x0F, 0x25, 0x40)
BORDER = RGBColor(0x23, 0x40, 0x60)
ACCENT = RGBColor(0x5E, 0xC4, 0xFF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT   = RGBColor(0xC0, 0xD4, 0xE8)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
MUTED  = RGBColor(0x70, 0x95, 0xB5)
FONT   = "Meiryo"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, x, y, w, h, text, size=14, color=TEXT, bold=False,
        align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb


def add_table_slide(title, subtitle, headers, rows, col_widths):
    s = prs.slide_layouts[6]
    slide = prs.slides.add_slide(s)
    bg(slide)
    box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
        title, size=24, color=WHITE, bold=True)
    if subtitle:
        box(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4),
            subtitle, size=12, color=MUTED)
    nrows = len(rows) + 1
    ncols = len(headers)
    total_w = sum(col_widths)
    left = Inches(0.5)
    top = Inches(1.45)
    height = Inches(min(5.6, 0.42 * nrows))
    gtbl = slide.shapes.add_table(nrows, ncols, left, top, Emu(int(total_w)), height)
    tbl = gtbl.table
    # disable banded style → manual colors
    tbl.first_row = False
    tbl.horz_banding = False
    for j, w in enumerate(col_widths):
        tbl.columns[j].width = Emu(int(w))
    # header
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = HEADBG
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = c.margin_right = Pt(4)
        c.margin_top = c.margin_bottom = Pt(2)
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = h
        r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = ACCENT; r.font.name = FONT
    # body
    for i, rowvals in enumerate(rows, start=1):
        for j, val in enumerate(rowvals):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = SURF
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Pt(4)
            c.margin_top = c.margin_bottom = Pt(1)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j >= 2 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(10)
            r.font.name = "Consolas" if j == 0 else FONT
            r.font.color.rgb = YELLOW if j == 2 else TEXT
    return slide


# ── Title slide ──────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
box(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.2),
    "Pitch Log — 評価基準カタログ", size=40, color=WHITE, bold=True)
box(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.8),
    "xT 以外の GSA 説明変数候補（全10カテゴリ）", size=20, color=ACCENT)
box(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.6),
    "凡例：★=GSA寄与度（多いほど強力） / 🟢簡単 🟡中 🔴大", size=14, color=MUTED)
box(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5),
    "出典: docs/feature_catalog.md", size=11, color=MUTED)

H4 = ["カラム名", "内容", "GSA価値", "実装"]
CW = [Emu(Inches(3.6)), Emu(Inches(5.3)), Emu(Inches(1.6)), Emu(Inches(1.3))]

# ── ① 動き・速度 ──
add_table_slide("① 動き・速度系", "🟢 簡単で高ROI（22選手分で約100カラム）", H4, [
    ["Ball_vx, Ball_vy", "ボール速度成分 (m/s)", "★★★", "🟢"],
    ["Ball_speed", "ボール速度 (m/s)", "★★★", "🟢"],
    ["Ball_acceleration", "ボール加速度 (m/s²)", "★★", "🟢"],
    ["Ball_direction", "ボール進行方向 (角度)", "★★", "🟢"],
    ["{team}_P{n}_acceleration", "各選手の加速度", "★★★", "🟢"],
    ["{team}_P{n}_vx, vy", "各選手の速度成分", "★★", "🟢"],
    ["{team}_P{n}_is_sprinting", "スプリント判定 (>5.5m/s)", "★★", "🟢"],
    ["{team}_P{n}_dist_traveled", "累積走行距離", "★", "🟢"],
], CW)

# ── ② 距離・位置関係 ──
add_table_slide("② 距離・位置関係", "🟢〜🟡", H4, [
    ["{team}_P{n}_dist_goal", "各選手からゴールまでの距離", "★★★", "🟢"],
    ["{team}_P{n}_dist_nearest_opp", "最近接相手選手との距離", "★★★", "🟡"],
    ["{team}_P{n}_n_opp_within_3m", "3m以内の相手数（プレッシャー）", "★★★", "🟡"],
    ["{team}_P{n}_n_opp_within_5m", "5m以内の相手選手数", "★★", "🟡"],
    ["Ball_dist_goal_home/away", "ボール→各ゴール距離", "★★", "🟢"],
    ["Ball_carrier_id", "現在のボール保持者ID", "★★★", "🟡"],
    ["Ball_carrier_team", "保持チーム（Home/Away/None）", "★★★", "🟡"],
    ["Pressure_on_carrier", "保持者への3m以内相手数", "★★★", "🟡"],
], CW)

# ── ③ チーム形状・組織 ──
add_table_slide("③ チーム形状・組織", "🟡 戦術的崩れを捉える", H4, [
    ["Home/Away_centroid_X,_Y", "チーム重心座標", "★★", "🟢"],
    ["Home/Away_width", "横幅 (Ymax - Ymin)", "★★", "🟢"],
    ["Home/Away_depth", "縦幅 (Xmax - Xmin)", "★★", "🟢"],
    ["Home/Away_compactness", "選手間平均距離", "★★", "🟡"],
    ["Home/Away_hull_area", "布陣カバー面積 (凸包)", "★★", "🟡"],
    ["Home/Away_last_def_line_X", "最終DFラインのX座標", "★★★", "🟢"],
    ["..._interline_dist_DF_MF", "DF-MF間距離", "★★", "🟡"],
    ["..._interline_dist_MF_FW", "MF-FW間距離", "★★", "🟡"],
], CW)

# ── ④ ゾーン拡張 ──
add_table_slide("④ ゾーン拡張", "🟢", H4, [
    ["Ball_Zone_5div", "縦5分割 (1〜5)", "★★", "🟢"],
    ["Ball_Zone_lateral", "横3分割 (left/center/right)", "★★", "🟢"],
    ["Ball_Grid20", "5×4=20領域グリッド", "★★", "🟢"],
    ["Ball_in_penalty_area", "ペナルティエリア内フラグ", "★★★", "🟢"],
    ["Ball_in_box_18yard", "18ヤードボックス内フラグ", "★★★", "🟢"],
    ["Ball_in_half_space", "ハーフスペース内フラグ", "★★", "🟢"],
    ["Ball_zone_transition", "ゾーン変化フラグ (進入/退出)", "★★★", "🟢"],
    ["Ball_attacking_third_dwell", "アタッキングサード滞在累計秒", "★★", "🟢"],
], CW)

# ── ⑤ ボール所有・フェーズ ──
add_table_slide("⑤ ボール所有・フェーズ", "🟡〜🔴", H4, [
    ["Ball_possession", "Home/Away/None", "★★★", "🟡"],
    ["Possession_duration_sec", "連続保持時間", "★★★", "🟡"],
    ["Possession_team_phase", "ビルドアップ/トランジ等", "★★", "🔴"],
    ["Defensive_phase", "ブロック構築/プレッシング等", "★★", "🔴"],
    ["Is_open_play", "オープンプレー中か", "★", "🟢"],
    ["Is_set_piece", "セットプレー中か", "★", "🟢"],
], CW)

# ── ⑥ パス・前進系 ──
add_table_slide("⑥ パス・前進系", "🟡〜🔴", H4, [
    ["Pass_event_flag", "パス発生フレーム", "★★★", "🟡"],
    ["Pass_xT_gain", "パスによるxT増加", "★★★", "🟡"],
    ["Pass_distance_m", "パス距離 (m)", "★★", "🟡"],
    ["Pass_forward_component_m", "パスの前進成分 (m)", "★★", "🟡"],
    ["Pass_breaks_line_count", "突破ライン数 (0〜3)", "★★★", "🔴"],
    ["Ball_forward_velocity", "ボール前進速度", "★★", "🟢"],
], CW)

# ── ⑦ 個人xT 変化系 ──
add_table_slide("⑦ 個人xT 変化系", "🟢（22人×4列で約88カラム）", H4, [
    ["{team}_P{n}_Delta_xT", "各選手のxT変化量", "★★★", "🟢"],
    ["{team}_P{n}_Cumulative_xT_gain", "各選手の累積xT獲得", "★★", "🟢"],
    ["{team}_P{n}_xT_smoothed_3f", "3フレーム移動平均xT", "★", "🟢"],
    ["{team}_P{n}_peak_xT_in_scene", "シーン中の最大xT", "★", "🟢"],
    ["{team}_P{n}_xT_acceleration", "xT変化の加速度", "★", "🟢"],
    ["{team}_xT_weighted_centroid_X", "xT加重重心", "★★", "🟡"],
], CW)

# ── ⑧ 時間ラグ ──
add_table_slide("⑧ 時間ラグ（GSA因果分析の本命）", "🟢 因果方向を特定する核心", H4, [
    ["Ball_xT_lag1", "1フレーム前 (0.1秒前)", "★★★", "🟢"],
    ["Ball_xT_lag5", "5フレーム前 (0.5秒前)", "★★★", "🟢"],
    ["Ball_xT_lag10", "10フレーム前 (1秒前)", "★★★", "🟢"],
    ["Ball_xT_lag20", "20フレーム前 (2秒前)", "★★", "🟢"],
    ["{team}_P{n}_xT_lag5", "各選手xTの0.5秒前", "★★★", "🟢"],
    ["{team}_P{n}_xT_lag10", "各選手xTの1秒前", "★★", "🟢"],
    ["{team}_P{n}_speed_lag5", "各選手速度の0.5秒前", "★★", "🟢"],
], CW)

# ── ⑨ dynamic_events 連動 ──
add_table_slide("⑨ dynamic_events.csv 連動", "🔴 外部データ統合", H4, [
    ["Is_pass_moment", "パス発生フレームフラグ", "★★★", "🔴"],
    ["Is_shot_moment", "シュート発生フレームフラグ", "★★★", "🔴"],
    ["Is_engagement_moment", "守備接触発生フラグ", "★★★", "🔴"],
    ["Is_off_ball_run_active", "オフボールラン進行中", "★★", "🔴"],
    ["Frame_xshot_start/end", "xshot値 (時系列)", "★★★", "🔴"],
    ["Cumulative_VAEP_attack/defend", "累積VAEP（攻撃/守備）", "★★", "🔴"],
    ["Pressing_chain_active", "プレッシングチェーン中", "★★", "🔴"],
], CW)

# ── ⑩ 目的変数候補 ──
add_table_slide("⑩ 目的変数候補の追加", "🟢〜🟡", H4, [
    ["Ball_xT_smoothed_5f", "ノイズ除去版 Ball_xT", "目的変数", "🟢"],
    ["Delta_Ball_xT_smoothed_3f", "ノイズ除去版 ΔxT", "目的変数", "🟢"],
    ["Ball_xT_max_in_5sec_window", "5秒ウィンドウ最大xT", "目的変数", "🟢"],
    ["Time_to_peak_xT_sec", "xTピークまでの秒数", "目的変数", "🟢"],
    ["Time_to_goal_sec", "ゴールまでの残り秒数", "目的変数", "🟢"],
    ["Is_attacking_third_entry", "アタッキングサード進入", "目的変数", "🟢"],
    ["Is_penalty_area_entry", "ペナルティエリア進入", "目的変数", "🟢"],
], CW)

# ── ロードマップ ──
add_table_slide("📊 推奨実装ロードマップ", "段階的に追加していく計画",
    ["Phase", "内容", "追加列", "時間"], [
    ["Phase 1", "動き・距離・ラグ（全🟢）", "+110", "30分"],
    ["Phase 2", "チーム形状・ライン（🟡）", "+50", "1〜2h"],
    ["Phase 3", "events統合（🔴）", "+20", "半日"],
], [Emu(Inches(2.0)), Emu(Inches(7.0)), Emu(Inches(1.6)), Emu(Inches(1.6))])

# ── 多重共線性 ──
s = prs.slides.add_slide(BLANK); bg(s)
box(s, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6),
    "⚠️ 多重共線性の注意", size=24, color=WHITE, bold=True)
box(s, Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.0),
    "以下のペアは両方入れない（GSA投入時はどちらか片方）：\n\n"
    "・位置(X,Y) と 個人xT  … xTは位置から計算される\n"
    "・個人xT と 個人GridID … GridIDも位置の関数\n"
    "・MAX_xT と 個人xT全員 … MAXは個人xTの集計\n"
    "・SUM_xT と 個人xT全員 … SUMは個人xTの集計\n"
    "・Ball_speed と Ball_vx/vy … speedは成分から計算\n"
    "・Pressure_on_carrier と n_opp_within_3m(carrier) … 同じ情報",
    size=16, color=TEXT)

out = "feature_catalog.pptx"
import sys
if len(sys.argv) > 1:
    out = sys.argv[1]
prs.save(out)
print("saved ->", out, "slides:", len(prs.slides._sldIdLst))
